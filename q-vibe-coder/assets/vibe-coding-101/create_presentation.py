from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Create presentation with widescreen dimensions
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
DARK_BG = RGBColor(15, 20, 25)  # Dark background
BLUE_ACCENT = RGBColor(29, 155, 240)  # Twitter blue
WHITE = RGBColor(231, 233, 234)
GRAY = RGBColor(113, 118, 123)
GREEN = RGBColor(0, 186, 124)
RED = RGBColor(244, 33, 46)

def add_dark_background(slide):
    """Add dark background to slide"""
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()
    # Send to back
    spTree = slide.shapes._spTree
    sp = background._element
    spTree.remove(sp)
    spTree.insert(2, sp)

def add_title_slide(title, subtitle=""):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    add_dark_background(slide)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        # Subtitle
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
        tf = sub_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(28)
        p.font.color.rgb = BLUE_ACCENT
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(title, bullets, quote=""):
    """Add a content slide with title and bullets"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    add_dark_background(slide)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Bullets
    bullet_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11.333), Inches(4.5))
    tf = bullet_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(24)
        p.font.color.rgb = WHITE
        p.space_before = Pt(12)
        p.space_after = Pt(6)

    # Quote at bottom
    if quote:
        quote_box = slide.shapes.add_textbox(Inches(1), Inches(6.3), Inches(11.333), Inches(0.8))
        tf = quote_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f'"{quote}"'
        p.font.size = Pt(20)
        p.font.italic = True
        p.font.color.rgb = BLUE_ACCENT
        p.alignment = PP_ALIGN.CENTER

    return slide

# SLIDE 1: Title
slide1 = add_title_slide("VIBE CODING 101", "Build and Deploy Real Websites with AI")
# Add instructor
inst_box = slide1.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.333), Inches(0.5))
tf = inst_box.text_frame
p = tf.paragraphs[0]
p.text = "Instructor: Guy Rymberg"
p.font.size = Pt(20)
p.font.color.rgb = GRAY
p.alignment = PP_ALIGN.CENTER

# SLIDE 2: What is Vibe Coding?
add_content_slide(
    "What is Vibe Coding?",
    [
        "Directing AI strategically through a proven process",
        "to build real software - not random trial and error",
        "",
        "• AI is your development partner",
        "• You bring vision and taste",
        "• AI brings engineering expertise",
        "• Together: real software, deployed"
    ]
)

# SLIDE 3: Common Mistakes
add_content_slide(
    "Common Vibe Coding Mistakes",
    [
        "❌ Random prompting with no plan",
        "❌ No architectural decisions upfront",
        "❌ Building without testing",
        "❌ Projects stuck on localhost forever",
        "❌ Losing context across sessions",
        "❌ Scope creep and endless tweaking"
    ],
    "Result: Chaos, frustration, abandoned projects"
)

# SLIDE 4: The Taste Principle
add_content_slide(
    "The Taste Principle",
    [
        "Non-coders control outcomes through:",
        "",
        "✓ Being specific about what you want",
        "✓ Bringing examples of things you like",
        "✓ Asking for options, then choosing",
        "✓ Recognizing quality when you see it"
    ],
    "Vibe coding is mostly about taste"
)

# SLIDE 5: 6-Phase Methodology
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_background(slide5)

title_box = slide5.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "The 6-Phase Vibe Coding Methodology"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Phase boxes - Row 1
phases_row1 = [("1", "VISION"), ("2", "CONSTRAINTS"), ("3", "ARCHITECTURE")]
for i, (num, name) in enumerate(phases_row1):
    left = Inches(1.5 + i * 3.5)
    box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.2), Inches(3), Inches(1.2))
    box.fill.solid()
    box.fill.fore_color.rgb = BLUE_ACCENT
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"{num}. {name}"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    tf.anchor = MSO_ANCHOR.MIDDLE

# Arrow
arrow_box = slide5.shapes.add_textbox(Inches(6), Inches(3.6), Inches(1), Inches(0.5))
tf = arrow_box.text_frame
p = tf.paragraphs[0]
p.text = "↓"
p.font.size = Pt(36)
p.font.color.rgb = BLUE_ACCENT
p.alignment = PP_ALIGN.CENTER

# Phase boxes - Row 2
phases_row2 = [("4", "BUILDING"), ("5", "TESTING"), ("6", "DEPLOYMENT")]
for i, (num, name) in enumerate(phases_row2):
    left = Inches(1.5 + i * 3.5)
    box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(4.2), Inches(3), Inches(1.2))
    box.fill.solid()
    box.fill.fore_color.rgb = BLUE_ACCENT
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"{num}. {name}"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    tf.anchor = MSO_ANCHOR.MIDDLE

# Quote
quote_box = slide5.shapes.add_textbox(Inches(1), Inches(6), Inches(11.333), Inches(0.8))
tf = quote_box.text_frame
p = tf.paragraphs[0]
p.text = '"Structure turns vibe coding into a reliable skill"'
p.font.size = Pt(20)
p.font.italic = True
p.font.color.rgb = BLUE_ACCENT
p.alignment = PP_ALIGN.CENTER

# SLIDE 6: Phase 1 - Vision
add_content_slide(
    "Phase 1: Vision",
    [
        "Ask yourself:",
        "",
        "• What are you building?",
        "• What problem does it solve?",
        "• Who is it for?",
        "• What does success look like?",
        "• What inspired you? (examples, references)"
    ],
    "If you can't explain it, you can't build it"
)

# SLIDE 7: Phase 2 - Constraints
add_content_slide(
    "Phase 2: Constraints",
    [
        "Define your boundaries:",
        "",
        "• Where will it run? (Web? Desktop? Mobile?)",
        "• What's the complexity? (Simple? Multi-session?)",
        "• What must it integrate with? (APIs? Services?)",
        "• What is NOT included? (Scope boundary)"
    ],
    "Constraints prevent scope creep"
)

# SLIDE 8: Phase 3 - Architecture
add_content_slide(
    "Phase 3: Architecture",
    [
        "Plan before building:",
        "",
        "• Break project into pieces (components)",
        "• Choose your stack (Next.js, Tailwind, etc.)",
        "• Design data flow (what talks to what)",
        "• Plan build sequence (what to build first)"
    ],
    "Decide WHERE it runs before HOW it works"
)

# SLIDE 9: Phase 4 - Building
add_content_slide(
    "Phase 4: Building",
    [
        "Iterative development:",
        "",
        "• Build in small, testable pieces",
        "• Test as you go (does it work?)",
        "• Commit regularly (save your work)",
        "• Stay within scope (resist feature creep)",
        "",
        'Watch for: "Oh, and it should also..."',
        'Response: "Great idea - let\'s note that for v2"'
    ]
)

# SLIDE 10: Phase 5 - Testing
add_content_slide(
    "Phase 5: Testing",
    [
        "Validate before deploying:",
        "",
        "• Does it work? (Functional testing)",
        "• Does it handle errors? (Edge cases)",
        "• Does it look right? (Visual verification)",
        "• Would a user understand it? (UX check)"
    ],
    "Good enough is good enough - avoid endless tweaking"
)

# SLIDE 11: Phase 6 - Deployment
add_content_slide(
    "Phase 6: Deployment",
    [
        "Get it live:",
        "",
        "• Deploy to production (Vercel)",
        "• Connect real services (domain, APIs)",
        "• Test the live URL",
        "• Share with the world"
    ],
    "A deployed app that works beats a perfect local prototype"
)

# SLIDE 12: Q-System
add_content_slide(
    "Managing Long Projects",
    [
        "Q-System Commands:",
        "",
        "/q-begin        → Start session, load context",
        "/q-end          → End session, save progress",
        "/q-checkpoint   → Save mid-session backup",
        "/q-status       → Check current state",
        "",
        "q-vibe-coder template:",
        "• vibe-coder-profile.md (your skills)",
        "• project.md (decisions & scope)",
        "• session-log.md (progress tracking)"
    ]
)

# SLIDE 13: Professional Tools
slide13 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_background(slide13)

title_box = slide13.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "Professional Tools"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# GitHub box
gh_box = slide13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(2), Inches(3.5), Inches(3))
gh_box.fill.solid()
gh_box.fill.fore_color.rgb = RGBColor(36, 41, 47)
gh_box.line.fill.background()

gh_text = slide13.shapes.add_textbox(Inches(2.2), Inches(2.3), Inches(3.1), Inches(2.5))
tf = gh_text.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "GitHub"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

p = tf.add_paragraph()
p.text = "\n• Version control"
p.font.size = Pt(18)
p.font.color.rgb = WHITE

p = tf.add_paragraph()
p.text = "• Code history"
p.font.size = Pt(18)
p.font.color.rgb = WHITE

p = tf.add_paragraph()
p.text = "• Collaboration"
p.font.size = Pt(18)
p.font.color.rgb = WHITE

# Arrow
arrow_box = slide13.shapes.add_textbox(Inches(5.8), Inches(3.2), Inches(1.5), Inches(0.8))
tf = arrow_box.text_frame
p = tf.paragraphs[0]
p.text = "push →"
p.font.size = Pt(24)
p.font.color.rgb = BLUE_ACCENT
p.alignment = PP_ALIGN.CENTER

# Vercel box
v_box = slide13.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.8), Inches(2), Inches(3.5), Inches(3))
v_box.fill.solid()
v_box.fill.fore_color.rgb = RGBColor(0, 0, 0)
v_box.line.color.rgb = WHITE

v_text = slide13.shapes.add_textbox(Inches(8), Inches(2.3), Inches(3.1), Inches(2.5))
tf = v_text.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Vercel"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

p = tf.add_paragraph()
p.text = "\n• Auto-deploy"
p.font.size = Pt(18)
p.font.color.rgb = WHITE

p = tf.add_paragraph()
p.text = "• Live URL"
p.font.size = Pt(18)
p.font.color.rgb = WHITE

p = tf.add_paragraph()
p.text = "• Preview builds"
p.font.size = Pt(18)
p.font.color.rgb = WHITE

# Quote
quote_box = slide13.shapes.add_textbox(Inches(1), Inches(5.8), Inches(11.333), Inches(0.8))
tf = quote_box.text_frame
p = tf.paragraphs[0]
p.text = '"Same tools the professionals use"'
p.font.size = Pt(20)
p.font.italic = True
p.font.color.rgb = BLUE_ACCENT
p.alignment = PP_ALIGN.CENTER

# SLIDE 14: What You'll Build
add_content_slide(
    "What You'll Build",
    [
        "By the end of this course:",
        "",
        "✓ A fully functional website deployed to production",
        "✓ Your own GitHub repository with clean code",
        "✓ Live Vercel URL you can share with anyone",
        "✓ q-vibe-coder project structure for future builds",
        "",
        "Course Format:",
        "• 2 live 1-on-1 sessions (90 min each)",
        "• Hands-on building during each session",
        "• Certificate of completion"
    ]
)

# SLIDE 15: Next Steps
slide15 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_background(slide15)

title_box = slide15.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "Ready to Start?"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

prereq_box = slide15.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11.333), Inches(2.5))
tf = prereq_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Prerequisites:"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = WHITE

prereqs = [
    "✓ Intro to Claude Code (or equivalent knowledge)",
    "✓ Claude Pro or Max subscription",
    "✓ GitHub account (free)",
    "✓ A project idea"
]
for prereq in prereqs:
    p = tf.add_paragraph()
    p.text = prereq
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    p.space_before = Pt(8)

# CTA Button
cta_box = slide15.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(5), Inches(5.333), Inches(1))
cta_box.fill.solid()
cta_box.fill.fore_color.rgb = BLUE_ACCENT
cta_box.line.fill.background()

cta_text = slide15.shapes.add_textbox(Inches(4), Inches(5.1), Inches(5.333), Inches(0.8))
tf = cta_text.text_frame
p = tf.paragraphs[0]
p.text = "ENROLL NOW - $249"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Contact
contact_box = slide15.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11.333), Inches(0.5))
tf = contact_box.text_frame
p = tf.paragraphs[0]
p.text = "Questions? Contact us at PeerLoop"
p.font.size = Pt(18)
p.font.color.rgb = GRAY
p.alignment = PP_ALIGN.CENTER

# Save presentation
output_path = r"C:\PeerLoop2\q-vibe-coder\assets\vibe-coding-101\Vibe_Coding_101_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
